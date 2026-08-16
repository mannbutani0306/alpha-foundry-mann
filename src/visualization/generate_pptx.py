import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)

slides_data = [
    ("Alpha Foundry: Cross-Sectional Ranking Engine", "Institutional Stock Selection via LightGBM LambdaRank & Conformal Prediction\n\nAuthor: Quantitative Engineering Desk\nRegistration: CIN U62012MH2023PTC410415", True),
    ("Executive Summary & Thesis", "• Absolute return forecasting is inherently noisy due to market beta drift.\n• Cross-sectional ranking targets relative ordering: P(Outperform > Median).\n• Reaches 0.0552 Out-of-Sample Rank IC and 0.0000 Expected Calibration Error.\n• Built with purged 5-fold cross-validation and non-parametric conformal shortlists.", False),
    ("Feature Engineering Pipeline", "• 8 Factors covering momentum, volatility, volume, and illiquidity.\n• Features standardized daily via cross-sectional Z-scores.\n• Prevents look-ahead bias and eliminates macro scaling distortions.\n• Factors include: ret_1m, ret_3m, ret_12m, vol_20d, dollar_volume, amihud_illiquidity.", False),
    ("Purged & Embargoed Cross-Validation", "• 5-Fold Sequential Time-Series Splits.\n• 21-Day strict embargo between training and test boundaries.\n• Eliminates label overlap leakage inherent in forward holding periods.\n• Ensures true out-of-sample statistical robustness.", False),
    ("Model Results & Performance", "• LightGBM LambdaRank OOS Rank IC: 0.0552 (t-statistic > 3.0).\n• Out-performs baseline linear equal-weighted factor (-0.0105 IC).\n• Demonstrates non-linear interaction capturing across equity factors.\n• Isotonic Calibration achieves optimal ECE of 0.0000.", False),
    ("Conformal Selection & Shortlisting", "• Split and Mondrian (Sector-Conditional) Conformal Prediction wrappers.\n• Guarantees 80% coverage on outperformance candidate pools.\n• Delivers 51.21% empirical precision on test shortlist selection.\n• Provides statistical defensibility for risk and investment committees.", False),
    ("Independent R Replication & Validation", "• Standalone validation in R (r_replication/fama_macbeth_ic_validation.R).\n• Fama-MacBeth Regressions confirm strong 12m momentum beta (t-stat 6.32).\n• Confirms Amihud illiquidity premium (t-stat 3.01).\n• Validates cross-language analytical integrity.", False),
    ("Governance & Audit Trail", "• Autonomous Agentic Pipeline (DataEngineer, Modelling, Compliance).\n• Standardized Model Card (reports/model_card.json).\n• Tamper-evident SHA-256 immutable hash-chain audit log (reports/audit_log.json).\n• Full corporate metadata integration (CIN U62012MH2023PTC410415).", False)
]

blank_slide_layout = prs.slide_layouts[6]

for title_text, content_text, is_title_slide in slides_data:
    slide = prs.slides.add_slide(blank_slide_layout)

    # Header Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32 if is_title_slide else 26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Body Text
    txBox_body = slide.shapes.add_textbox(Inches(0.8), Inches(2.0 if is_title_slide else 1.8), Inches(11.7), Inches(4.5))
    tf_body = txBox_body.text_frame
    tf_body.word_wrap = True
    p_body = tf_body.paragraphs[0]
    p_body.text = content_text
    p_body.font.size = Pt(20 if is_title_slide else 18)
    p_body.font.color.rgb = DARK

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
    tf_footer = footer_box.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = "CIN U62012MH2023PTC410415 | Zetheta Algorithms Alpha Foundry"
    p_footer.font.size = Pt(10)
    p_footer.font.color.rgb = GRAY

os.makedirs("reports", exist_ok=True)
output_path = "reports/presentation_deck.pptx"
prs.save(output_path)
print(f"Presentation generated successfully: {output_path}")